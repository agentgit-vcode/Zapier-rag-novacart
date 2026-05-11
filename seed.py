"""Seed Pinecone with chunked NovaCart documents.

Reads the documents from Files/, chunks them, embeds via OpenAI,
and upserts into the Pinecone index.

Usage:
    python seed.py
"""

import os
import sys

from dotenv import load_dotenv

load_dotenv()

from vectorstore import upsert_chunks  # noqa: E402

# --- Document Chunking ---

def chunk_text(text: str, chunk_size: int = 500, overlap: int = 100) -> list[str]:
    """Split text into overlapping chunks by character count."""
    chunks = []
    start = 0
    while start < len(text):
        end = start + chunk_size
        chunks.append(text[start:end])
        start += chunk_size - overlap
    return [c.strip() for c in chunks if c.strip()]


def load_docx(file_path: str) -> str:
    """Extract text from a .docx file."""
    try:
        import docx
    except ImportError:
        print("ERROR: python-docx is required. Install it: pip install python-docx")
        sys.exit(1)

    doc = docx.Document(file_path)
    return "\n".join(para.text for para in doc.paragraphs if para.text.strip())


def load_pptx(file_path: str) -> str:
    """Extract text from a .pptx file."""
    try:
        from pptx import Presentation
    except ImportError:
        print("ERROR: python-pptx is required. Install it: pip install python-pptx")
        sys.exit(1)

    prs = Presentation(file_path)
    text_parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        text_parts.append(para.text)
    return "\n".join(text_parts)


def load_pdf(file_path: str) -> str:
    """Extract text from a .pdf file."""
    try:
        import PyPDF2
    except ImportError:
        print("ERROR: PyPDF2 is required. Install it: pip install PyPDF2")
        sys.exit(1)

    with open(file_path, "rb") as f:
        reader = PyPDF2.PdfReader(f)
        text_parts = []
        for page in reader.pages:
            text = page.extract_text()
            if text:
                text_parts.append(text)
    return "\n".join(text_parts)


def load_xlsx(file_path: str) -> str:
    """Extract text from a .xlsx file by converting rows to readable strings."""
    try:
        import openpyxl
    except ImportError:
        print("ERROR: openpyxl is required. Install it: pip install openpyxl")
        sys.exit(1)

    wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
    text_parts = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = list(ws.iter_rows(values_only=True))
        if not rows:
            continue

        headers = [str(h) if h else f"col_{i}" for i, h in enumerate(rows[0])]
        text_parts.append(f"Sheet: {sheet_name}")
        text_parts.append(f"Columns: {', '.join(headers)}")

        for row in rows[1:]:
            row_text = " | ".join(
                f"{headers[i]}: {val}" for i, val in enumerate(row) if val is not None
            )
            if row_text.strip():
                text_parts.append(row_text)

    wb.close()
    return "\n".join(text_parts)


# --- Document Registry ---

DOCUMENTS = [
    {
        "file_name": "NovaCart Company Annual Report 2024.docx",
        "file_type": "annual_report",
        "loader": load_docx,
    },
    {
        "file_name": "NovaCart Company Intro.pptx",
        "file_type": "presentation",
        "loader": load_pptx,
    },
    {
        "file_name": "NovaCart_Product_Catalog.pdf",
        "file_type": "product_catalog",
        "loader": load_pdf,
    },
    {
        "file_name": "SKU_Weekly_Sales_Conversion_3Y_with_Revenue.xlsx",
        "file_type": "sales_data",
        "loader": load_xlsx,
    },
]


def seed():
    """Load all documents, chunk them, and upsert into Pinecone."""
    all_chunks = []

    for doc in DOCUMENTS:
        file_path = os.path.join("Files", doc["file_name"])

        if not os.path.exists(file_path):
            print(f"SKIP: {file_path} not found")
            continue

        print(f"Loading: {doc['file_name']}...")
        text = doc["loader"](file_path)

        if not text.strip():
            print(f"  WARNING: No text extracted from {doc['file_name']}")
            continue

        chunks = chunk_text(text)
        print(f"  Extracted {len(chunks)} chunks ({len(text)} chars)")

        for i, chunk in enumerate(chunks):
            all_chunks.append({
                "chunk_id": f"{doc['file_type']}_{i:04d}",
                "text": chunk,
                "metadata": {
                    "file_name": doc["file_name"],
                    "file_type": doc["file_type"],
                    "chunk_index": i,
                },
            })

    if not all_chunks:
        print("ERROR: No chunks to upsert. Check that Files/ directory contains documents.")
        sys.exit(1)

    print(f"\nUpserting {len(all_chunks)} chunks to Pinecone index '{os.getenv('PINECONE_INDEX_NAME', 'novacart-rag')}'...")
    count = upsert_chunks(all_chunks)
    print(f"Done! {count} chunks indexed.")


if __name__ == "__main__":
    seed()
